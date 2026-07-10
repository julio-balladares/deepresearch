from __future__ import annotations

import csv
import io
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup

from .cache import LocalCache, PAGE_CACHE_TTL_SECONDS, PAGES_CACHE_BUCKET
from .config import DEFAULT_TIMEOUT, logger
from .models import EvidenceSource
from .utils import compact_spaces, domain_from_url, truncate_text


MAX_DOCUMENT_PAGES = 60
MAX_SPREADSHEET_ROWS = 300
MAX_DOWNLOAD_BYTES = 8 * 1024 * 1024

TEXT_SOURCE_TYPES = {"csv", "json", "text", "markdown", "xml"}

BLOCKED_CONTENT_MARKERS = (
    "client challenge",
    "tiktok - make your day",
    "access denied",
    "just a moment",
    "enable javascript",
    "verify you are human",
    "checking your browser",
    "cloudflare ray id",
)


class SourceReader:
    def __init__(
        self,
        session: requests.Session,
        cache: LocalCache,
        max_chars_per_source: int,
    ) -> None:
        self.session = session
        self.cache = cache
        self.max_chars_per_source = max(500, int(max_chars_per_source))

    def fetch_page_text(self, url: str) -> tuple[str, bool, str]:
        clean_url = compact_spaces(url)

        if not self._is_supported_url(clean_url):
            return "", False, "invalid"

        cached = self.cache.get(
            PAGES_CACHE_BUCKET,
            clean_url,
            PAGE_CACHE_TTL_SECONDS,
        )

        if isinstance(cached, dict):
            return (
                compact_spaces(cached.get("text", "")),
                bool(cached.get("fetched")),
                compact_spaces(cached.get("source_type", "web")) or "web",
            )

        source_type = "web"
        text = ""

        try:
            response = self.session.get(
                clean_url,
                timeout=DEFAULT_TIMEOUT,
                allow_redirects=True,
                stream=True,
            )
            response.raise_for_status()

            content_type = response.headers.get("Content-Type", "").lower()
            final_url = response.url or clean_url
            source_type = self._detect_source_type(final_url, content_type)

            content = self._read_limited_content(response, MAX_DOWNLOAD_BYTES)

            if source_type == "html":
                text = self._extract_html_text(self._decode_text(content))
            elif source_type == "pdf":
                text = self._extract_pdf_text(content)
            elif source_type == "docx":
                text = self._extract_docx_text(content)
            elif source_type == "pptx":
                text = self._extract_pptx_text(content)
            elif source_type == "xlsx":
                text = self._extract_xlsx_text(content)
            elif source_type == "csv":
                text = self._extract_csv_text(content)
            elif source_type in TEXT_SOURCE_TYPES:
                text = self._decode_text(content)
            else:
                text = ""

        except requests.exceptions.RequestException as exc:
            logger.debug("Could not fetch %s: %s", clean_url, exc)
            return "", False, source_type
        except Exception as exc:
            logger.debug("Could not read %s: %s", clean_url, exc)
            return "", False, source_type

        text = truncate_text(compact_spaces(text), self.max_chars_per_source)

        if not self._is_page_content_usable(text):
            return "", False, source_type

        payload = {
            "text": text,
            "fetched": True,
            "source_type": source_type,
        }

        self.cache.set(PAGES_CACHE_BUCKET, clean_url, payload)

        return text, True, source_type

    def build_evidence(
        self,
        search_item: dict[str, str],
        query: str,
        subquestion: str,
        fetch_pages: bool = True,
    ) -> EvidenceSource:
        url = compact_spaces(search_item.get("url", ""))
        title = compact_spaces(search_item.get("title", ""))
        snippet = compact_spaces(search_item.get("snippet", ""))

        text = ""
        fetched = False
        source_type = "web"

        if fetch_pages and url:
            text, fetched, source_type = self.fetch_page_text(url)

        return EvidenceSource(
            title=title,
            url=url,
            domain=domain_from_url(url)
            or compact_spaces(search_item.get("domain", "")),
            snippet=snippet,
            extracted_text=text,
            query=compact_spaces(query),
            subquestion=compact_spaces(subquestion),
            fetched=fetched,
            source_type=source_type,
        )

    @staticmethod
    def _read_limited_content(
        response: requests.Response,
        max_bytes: int,
    ) -> bytes:
        chunks: list[bytes] = []
        total = 0

        for chunk in response.iter_content(chunk_size=64 * 1024):
            if not chunk:
                continue

            total += len(chunk)

            if total > max_bytes:
                raise ValueError(f"Downloaded content exceeds {max_bytes} bytes.")

            chunks.append(chunk)

        return b"".join(chunks)

    @staticmethod
    def _is_supported_url(url: str) -> bool:
        parsed = urlparse(url)

        return parsed.scheme in {"http", "https"} and bool(parsed.netloc)

    @staticmethod
    def _is_page_content_usable(text: str) -> bool:
        normalized = compact_spaces(text).lower()

        if not normalized:
            return False

        return not any(marker in normalized for marker in BLOCKED_CONTENT_MARKERS)

    @staticmethod
    def _detect_source_type(url: str, content_type: str) -> str:
        path = urlparse(url).path.lower()
        normalized_content_type = (content_type or "").split(";")[0].strip().lower()

        extension_map = {
            ".pdf": "pdf",
            ".docx": "docx",
            ".pptx": "pptx",
            ".xlsx": "xlsx",
            ".xlsm": "xlsx",
            ".csv": "csv",
            ".txt": "text",
            ".md": "markdown",
            ".json": "json",
            ".xml": "xml",
        }

        for extension, source_type in extension_map.items():
            if path.endswith(extension):
                return source_type

        content_type_map = {
            "application/pdf": "pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
            "application/vnd.ms-excel.sheet.macroenabled.12": "xlsx",
            "text/csv": "csv",
            "application/csv": "csv",
            "text/plain": "text",
            "text/markdown": "markdown",
            "application/json": "json",
            "application/xml": "xml",
            "text/xml": "xml",
            "text/html": "html",
            "application/xhtml+xml": "html",
        }

        if normalized_content_type in content_type_map:
            return content_type_map[normalized_content_type]

        if not normalized_content_type:
            return "html"

        if normalized_content_type.startswith("text/"):
            return "text"

        return "other"

    @staticmethod
    def _extract_html_text(html: str) -> str:
        soup = BeautifulSoup(html or "", "html.parser")

        for tag in soup(
            [
                "script",
                "style",
                "noscript",
                "svg",
                "img",
                "iframe",
                "header",
                "footer",
                "nav",
                "aside",
                "form",
                "button",
            ]
        ):
            tag.decompose()

        title = soup.title.get_text(" ", strip=True) if soup.title else ""
        main_text = soup.get_text(separator=" ", strip=True)

        return compact_spaces(f"{title}. {main_text}")

    @staticmethod
    def _extract_pdf_text(content: bytes) -> str:
        try:
            from pypdf import PdfReader
        except ImportError:
            logger.warning("Cannot read PDF: pypdf is not installed.")
            return ""

        try:
            reader = PdfReader(io.BytesIO(content))
            pages: list[str] = []

            for page in reader.pages[:MAX_DOCUMENT_PAGES]:
                pages.append(page.extract_text() or "")

            return compact_spaces("\n".join(pages))
        except Exception as exc:
            logger.debug("Could not extract PDF text: %s", exc)
            return ""

    @staticmethod
    def _extract_docx_text(content: bytes) -> str:
        try:
            from docx import Document
        except ImportError:
            logger.warning("Cannot read DOCX: python-docx is not installed.")
            return ""

        try:
            document = Document(io.BytesIO(content))
            parts: list[str] = []

            for paragraph in document.paragraphs:
                if paragraph.text:
                    parts.append(paragraph.text)

            for table in document.tables:
                for row in table.rows:
                    values = [cell.text for cell in row.cells if cell.text]
                    if values:
                        parts.append(" | ".join(values))

            return compact_spaces("\n".join(parts))
        except Exception as exc:
            logger.debug("Could not extract DOCX text: %s", exc)
            return ""

    @staticmethod
    def _extract_pptx_text(content: bytes) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            logger.warning("Cannot read PPTX: python-pptx is not installed.")
            return ""

        try:
            presentation = Presentation(io.BytesIO(content))
            parts: list[str] = []

            for slide_number, slide in enumerate(presentation.slides, start=1):
                parts.append(f"Slide {slide_number}")

                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        parts.append(str(text))

            return compact_spaces("\n".join(parts))
        except Exception as exc:
            logger.debug("Could not extract PPTX text: %s", exc)
            return ""

    @staticmethod
    def _extract_xlsx_text(content: bytes) -> str:
        try:
            from openpyxl import load_workbook
        except ImportError:
            logger.warning("Cannot read XLSX: openpyxl is not installed.")
            return ""

        workbook = None

        try:
            workbook = load_workbook(
                io.BytesIO(content),
                read_only=True,
                data_only=True,
            )

            parts: list[str] = []

            for sheet in workbook.worksheets:
                parts.append(f"Sheet: {sheet.title}")

                for row_index, row in enumerate(
                    sheet.iter_rows(values_only=True),
                    start=1,
                ):
                    if row_index > MAX_SPREADSHEET_ROWS:
                        parts.append("[TRUNCATED: too many rows]")
                        break

                    values = [
                        str(value)
                        for value in row
                        if value is not None and str(value).strip()
                    ]

                    if values:
                        parts.append(" | ".join(values))

            return compact_spaces("\n".join(parts))
        except Exception as exc:
            logger.debug("Could not extract XLSX text: %s", exc)
            return ""
        finally:
            if workbook is not None:
                workbook.close()

    @classmethod
    def _extract_csv_text(cls, content: bytes) -> str:
        text = cls._decode_text(content)

        if not text:
            return ""

        sample = text[:4096]

        try:
            dialect = csv.Sniffer().sniff(sample)
        except csv.Error:
            dialect = csv.excel

        rows: list[str] = []

        try:
            reader = csv.reader(io.StringIO(text), dialect)

            for row_index, row in enumerate(reader, start=1):
                if row_index > MAX_SPREADSHEET_ROWS:
                    rows.append("[TRUNCATED: too many rows]")
                    break

                values = [value.strip() for value in row if value and value.strip()]

                if values:
                    rows.append(" | ".join(values))

        except csv.Error:
            return text

        return compact_spaces("\n".join(rows))

    @staticmethod
    def _decode_text(content: bytes) -> str:
        for encoding in ("utf-8-sig", "utf-8", "latin-1"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue

        return content.decode("utf-8", errors="ignore")


__all__ = ["SourceReader"]
